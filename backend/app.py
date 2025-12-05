from flask import Flask, request, jsonify, send_file
import qrcode
import io
from flask_cors import CORS
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
from PyPDF2 import PdfReader, PdfWriter, PdfMerger
import pdfplumber
import os

# Chuyển đổi PDF sang Word
from pdf2docx import Converter
import tempfile

# Chuyển đổi PDF sang Excel
import tabula
import pandas as pd

# Rate limiting
from flask_limiter import Limiter
from flask_limiter.util import get_remote_address

app = Flask(__name__)
CORS(app)

# Cấu hình bảo mật
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024  # 50MB max file size

# Rate limiter
limiter = Limiter(
    app=app,
    key_func=get_remote_address,
    default_limits=["100 per hour", "20 per minute"],
    storage_uri="memory://"
)

# Allowed file extensions
ALLOWED_PDF_EXTENSIONS = {'pdf'}
ALLOWED_IMAGE_EXTENSIONS = {'png', 'jpg', 'jpeg', 'gif', 'bmp', 'webp'}

def allowed_file(filename, allowed_extensions):
    """Kiểm tra file extension có hợp lệ không"""
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in allowed_extensions

def validate_pdf_file(file):
    """Validate PDF file"""
    if not file:
        return False, "No file uploaded"
    if file.filename == '':
        return False, "Empty filename"
    if not allowed_file(file.filename, ALLOWED_PDF_EXTENSIONS):
        return False, "Only PDF files are allowed"
    return True, None

# Route cho trang chủ để xác nhận backend hoạt động
@app.route("/")
def home():
    return "ZapTools Backend is running!"

@app.route('/api/merge-pdfs', methods=['POST'])
@limiter.limit("10 per minute")
def merge_pdfs():
    try:
        files = request.files.getlist('pdfs')
        if not files or len(files) < 2:
            return jsonify({'error': 'Upload at least 2 PDF files'}), 400
        if len(files) > 20:
            return jsonify({'error': 'Maximum 20 files allowed'}), 400

        # Validate all files
        for f in files:
            is_valid, error_msg = validate_pdf_file(f)
            if not is_valid:
                return jsonify({'error': error_msg}), 400

        merger = PdfMerger()
        for f in files:
            merger.append(f)

        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as out_pdf:
            merger.write(out_pdf)
            out_pdf_path = out_pdf.name

        return send_file(out_pdf_path, as_attachment=True, download_name='merged.pdf', mimetype='application/pdf')
    except Exception as e:
        return jsonify({'error': f'Error merging PDFs: {str(e)}'}), 500

@app.route('/api/resize-pdf', methods=['POST'])
@limiter.limit("15 per minute")
def resize_pdf():
    try:
        if 'pdf' not in request.files:
            return jsonify({'error': 'No PDF file uploaded'}), 400

        pdf_file = request.files['pdf']
        is_valid, error_msg = validate_pdf_file(pdf_file)
        if not is_valid:
            return jsonify({'error': error_msg}), 400

        # Validate dimensions
        try:
            width = float(request.form.get('width', 595))  # A4 width default
            height = float(request.form.get('height', 842))  # A4 height default

            if width <= 0 or height <= 0 or width > 5000 or height > 5000:
                return jsonify({'error': 'Invalid dimensions (must be between 1-5000)'}), 400
        except ValueError:
            return jsonify({'error': 'Invalid dimension values'}), 400

        reader = PdfReader(pdf_file)
        writer = PdfWriter()
        for page in reader.pages:
            page.mediabox.upper_right = (width, height)
            writer.add_page(page)

        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as out_pdf:
            writer.write(out_pdf)
            out_pdf_path = out_pdf.name

        return send_file(out_pdf_path, as_attachment=True, download_name='resized.pdf', mimetype='application/pdf')
    except Exception as e:
        return jsonify({'error': f'Error resizing PDF: {str(e)}'}), 500

@app.route('/api/split-pdf', methods=['POST'])
@limiter.limit("15 per minute")
def split_pdf():
    try:
        if 'pdf' not in request.files:
            return jsonify({'error': 'No PDF file uploaded'}), 400

        pdf_file = request.files['pdf']
        is_valid, error_msg = validate_pdf_file(pdf_file)
        if not is_valid:
            return jsonify({'error': error_msg}), 400

        # Validate page numbers
        try:
            start = int(request.form.get('start', 1))
            end = int(request.form.get('end', 1))

            if start < 1 or end < 1:
                return jsonify({'error': 'Page numbers must be greater than 0'}), 400
            if start > end:
                return jsonify({'error': 'Start page must be less than or equal to end page'}), 400
        except ValueError:
            return jsonify({'error': 'Invalid page numbers'}), 400

        reader = PdfReader(pdf_file)
        total_pages = len(reader.pages)

        if end > total_pages:
            return jsonify({'error': f'End page exceeds total pages ({total_pages})'}), 400

        writer = PdfWriter()
        for i in range(start-1, end):
            writer.add_page(reader.pages[i])

        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as out_pdf:
            writer.write(out_pdf)
            out_pdf_path = out_pdf.name

        return send_file(out_pdf_path, as_attachment=True, download_name='split.pdf', mimetype='application/pdf')
    except Exception as e:
        return jsonify({'error': f'Error splitting PDF: {str(e)}'}), 500

@app.route('/api/pdf-to-powerpoint', methods=['POST'])
@limiter.limit("5 per minute")
def pdf_to_powerpoint():
    try:
        if 'pdf' not in request.files:
            return jsonify({'error': 'No PDF file uploaded'}), 400

        pdf_file = request.files['pdf']
        is_valid, error_msg = validate_pdf_file(pdf_file)
        if not is_valid:
            return jsonify({'error': error_msg}), 400

        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as temp_pdf:
            pdf_file.save(temp_pdf.name)

            # Check page count to prevent abuse
            reader = PdfReader(temp_pdf.name)
            if len(reader.pages) > 50:
                return jsonify({'error': 'PDF has too many pages (max 50)'}), 400

            images = convert_from_path(temp_pdf.name)
            pptx_path = temp_pdf.name.replace('.pdf', '.pptx')
            prs = Presentation()
            blank_slide_layout = prs.slide_layouts[6]

            for img in images:
                slide = prs.slides.add_slide(blank_slide_layout)
                img_temp = tempfile.NamedTemporaryFile(suffix='.png', delete=False)
                img.save(img_temp.name, 'PNG')
                slide.shapes.add_picture(img_temp.name, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)

            prs.save(pptx_path)

        return send_file(pptx_path, as_attachment=True, download_name='converted.pptx', mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')
    except Exception as e:
        return jsonify({'error': f'Error converting PDF to PowerPoint: {str(e)}'}), 500


@app.route('/api/pdf-to-excel', methods=['POST'])
@limiter.limit("5 per minute")
def pdf_to_excel():
    try:
        if 'pdf' not in request.files:
            return jsonify({'error': 'No PDF file uploaded'}), 400

        pdf_file = request.files['pdf']
        is_valid, error_msg = validate_pdf_file(pdf_file)
        if not is_valid:
            return jsonify({'error': error_msg}), 400

        # Lưu PDF tạm
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as temp_pdf:
            pdf_file.save(temp_pdf.name)

            # Check page count to prevent abuse
            reader = PdfReader(temp_pdf.name)
            if len(reader.pages) > 50:
                return jsonify({'error': 'PDF has too many pages (max 50)'}), 400

            tables = []  # lưu tất cả bảng tìm được

            # Mở PDF bằng pdfplumber
            with pdfplumber.open(temp_pdf.name) as pdf:
                for page_number, page in enumerate(pdf.pages, start=1):
                    extracted_tables = page.extract_tables()

                    if extracted_tables:
                        for tbl in extracted_tables:
                            # Bỏ qua bảng rỗng
                            if len(tbl) > 1:
                                df = pd.DataFrame(tbl[1:], columns=tbl[0])
                                tables.append(df)

            if not tables:
                return jsonify({'error': 'No tables found in PDF'}), 400

            # Xuất ra file Excel
            excel_path = temp_pdf.name.replace('.pdf', '.xlsx')

            with pd.ExcelWriter(excel_path, engine='openpyxl') as writer:
                for i, df in enumerate(tables):
                    df.to_excel(writer, sheet_name=f'Sheet{i+1}', index=False)

        return send_file(
            excel_path,
            as_attachment=True,
            download_name='converted.xlsx',
            mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
        )
    except Exception as e:
        return jsonify({'error': f'Error converting PDF to Excel: {str(e)}'}), 500

@app.route('/api/pdf-to-word', methods=['POST'])
@limiter.limit("5 per minute")
def pdf_to_word():
    try:
        if 'pdf' not in request.files:
            return jsonify({'error': 'No PDF file uploaded'}), 400

        pdf_file = request.files['pdf']
        is_valid, error_msg = validate_pdf_file(pdf_file)
        if not is_valid:
            return jsonify({'error': error_msg}), 400

        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as temp_pdf:
            pdf_file.save(temp_pdf.name)

            # Check page count to prevent abuse
            reader = PdfReader(temp_pdf.name)
            if len(reader.pages) > 50:
                return jsonify({'error': 'PDF has too many pages (max 50)'}), 400

            word_path = temp_pdf.name.replace('.pdf', '.docx')
            cv = Converter(temp_pdf.name)
            cv.convert(word_path, start=0, end=None)
            cv.close()

        return send_file(word_path, as_attachment=True, download_name='converted.docx', mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document')
    except Exception as e:
        return jsonify({'error': f'Error converting PDF to Word: {str(e)}'}), 500


@app.route('/api/pdf-to-image', methods=['POST'])
@limiter.limit("5 per minute")
def pdf_to_image():
    try:
        if 'pdf' not in request.files:
            return jsonify({'error': 'No PDF file uploaded'}), 400

        pdf_file = request.files['pdf']
        is_valid, error_msg = validate_pdf_file(pdf_file)
        if not is_valid:
            return jsonify({'error': error_msg}), 400

        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as temp_pdf:
            pdf_file.save(temp_pdf.name)

            # Check page count to prevent abuse
            reader = PdfReader(temp_pdf.name)
            if len(reader.pages) > 20:
                return jsonify({'error': 'PDF has too many pages (max 20)'}), 400

            # Convert PDF to images
            images = convert_from_path(temp_pdf.name, dpi=150)

            # Create a zip file with all images
            import zipfile
            zip_path = temp_pdf.name.replace('.pdf', '.zip')

            with zipfile.ZipFile(zip_path, 'w') as zipf:
                for i, img in enumerate(images):
                    img_temp = tempfile.NamedTemporaryFile(suffix='.png', delete=False)
                    img.save(img_temp.name, 'PNG')
                    zipf.write(img_temp.name, f'page_{i+1}.png')
                    os.unlink(img_temp.name)

            os.unlink(temp_pdf.name)

        return send_file(zip_path, as_attachment=True, download_name='images.zip', mimetype='application/zip')
    except Exception as e:
        return jsonify({'error': f'Error converting PDF to images: {str(e)}'}), 500

@app.route('/api/qrwifi', methods=['POST'])
@limiter.limit("20 per minute")
def generate_wifi_qr():
    try:
        data = request.json
        if not data:
            return jsonify({'error': 'No data provided'}), 400

        ssid = data.get('ssid', '').strip()
        password = data.get('password', '').strip()
        encryption = data.get('encryption', 'WPA').strip()

        # Validate input
        if not ssid:
            return jsonify({'error': 'SSID is required'}), 400
        if len(ssid) > 32:
            return jsonify({'error': 'SSID too long (max 32 characters)'}), 400
        if encryption not in ['WPA', 'WPA2', 'WEP', 'nopass']:
            return jsonify({'error': 'Invalid encryption type'}), 400
        if encryption != 'nopass' and len(password) > 63:
            return jsonify({'error': 'Password too long (max 63 characters)'}), 400

        qr_text = f"WIFI:T:{encryption};S:{ssid};P:{password};;"
        img = qrcode.make(qr_text)
        buf = io.BytesIO()
        img.save(buf, format='PNG')
        buf.seek(0)
        return send_file(buf, mimetype='image/png')
    except Exception as e:
        return jsonify({'error': f'Error generating QR code: {str(e)}'}), 500

@app.route('/api/image-compress', methods=['POST'])
@limiter.limit("15 per minute")
def compress_image():
    try:
        if 'image' not in request.files:
            return jsonify({'error': 'No image file uploaded'}), 400

        image_file = request.files['image']
        if not image_file or image_file.filename == '':
            return jsonify({'error': 'Empty filename'}), 400
        if not allowed_file(image_file.filename, ALLOWED_IMAGE_EXTENSIONS):
            return jsonify({'error': 'Invalid image format'}), 400

        # Get quality parameter (1-100)
        try:
            quality = int(request.form.get('quality', 85))
            if quality < 1 or quality > 100:
                quality = 85
        except ValueError:
            quality = 85

        from PIL import Image
        img = Image.open(image_file)

        # Convert RGBA to RGB if necessary
        if img.mode in ('RGBA', 'LA', 'P'):
            background = Image.new('RGB', img.size, (255, 255, 255))
            if img.mode == 'P':
                img = img.convert('RGBA')
            background.paste(img, mask=img.split()[-1] if img.mode in ('RGBA', 'LA') else None)
            img = background

        with tempfile.NamedTemporaryFile(suffix='.jpg', delete=False) as temp_img:
            img.save(temp_img.name, 'JPEG', quality=quality, optimize=True)
            temp_img_path = temp_img.name

        return send_file(temp_img_path, as_attachment=True, download_name='compressed.jpg', mimetype='image/jpeg')
    except Exception as e:
        return jsonify({'error': f'Error compressing image: {str(e)}'}), 500

@app.route('/api/image-resize', methods=['POST'])
@limiter.limit("15 per minute")
def resize_image():
    try:
        if 'image' not in request.files:
            return jsonify({'error': 'No image file uploaded'}), 400

        image_file = request.files['image']
        if not image_file or image_file.filename == '':
            return jsonify({'error': 'Empty filename'}), 400
        if not allowed_file(image_file.filename, ALLOWED_IMAGE_EXTENSIONS):
            return jsonify({'error': 'Invalid image format'}), 400

        # Get dimensions
        try:
            width = int(request.form.get('width', 800))
            height = int(request.form.get('height', 600))

            if width <= 0 or height <= 0 or width > 10000 or height > 10000:
                return jsonify({'error': 'Invalid dimensions (must be between 1-10000)'}), 400
        except ValueError:
            return jsonify({'error': 'Invalid dimension values'}), 400

        from PIL import Image
        img = Image.open(image_file)
        img = img.resize((width, height), Image.Resampling.LANCZOS)

        # Preserve original format
        original_ext = image_file.filename.rsplit('.', 1)[1].lower()
        img_format = original_ext.upper()
        if img_format == 'JPG':
            img_format = 'JPEG'

        with tempfile.NamedTemporaryFile(suffix=f'.{original_ext}', delete=False) as temp_img:
            img.save(temp_img.name, img_format)
            temp_img_path = temp_img.name

        return send_file(temp_img_path, as_attachment=True, download_name=f'resized.{original_ext}', mimetype=f'image/{original_ext}')
    except Exception as e:
        return jsonify({'error': f'Error resizing image: {str(e)}'}), 500

@app.route('/api/image-convert', methods=['POST'])
@limiter.limit("15 per minute")
def convert_image():
    try:
        if 'image' not in request.files:
            return jsonify({'error': 'No image file uploaded'}), 400

        image_file = request.files['image']
        if not image_file or image_file.filename == '':
            return jsonify({'error': 'Empty filename'}), 400
        if not allowed_file(image_file.filename, ALLOWED_IMAGE_EXTENSIONS):
            return jsonify({'error': 'Invalid image format'}), 400

        # Get target format
        target_format = request.form.get('format', 'png').lower()
        if target_format not in ['png', 'jpg', 'jpeg', 'webp', 'bmp']:
            return jsonify({'error': 'Invalid target format'}), 400

        from PIL import Image
        img = Image.open(image_file)

        # Convert format
        if target_format in ['jpg', 'jpeg']:
            if img.mode in ('RGBA', 'LA', 'P'):
                background = Image.new('RGB', img.size, (255, 255, 255))
                if img.mode == 'P':
                    img = img.convert('RGBA')
                background.paste(img, mask=img.split()[-1] if img.mode in ('RGBA', 'LA') else None)
                img = background
            img_format = 'JPEG'
            file_ext = 'jpg'
        else:
            img_format = target_format.upper()
            file_ext = target_format

        with tempfile.NamedTemporaryFile(suffix=f'.{file_ext}', delete=False) as temp_img:
            img.save(temp_img.name, img_format)
            temp_img_path = temp_img.name

        return send_file(temp_img_path, as_attachment=True, download_name=f'converted.{file_ext}', mimetype=f'image/{file_ext}')
    except Exception as e:
        return jsonify({'error': f'Error converting image: {str(e)}'}), 500

@app.route('/api/feedback', methods=['POST'])
@limiter.limit("10 per hour")
def submit_feedback():
    """Receive user feedback and log it"""
    try:
        data = request.get_json()

        if not data or 'message' not in data:
            return jsonify({'error': 'Message is required'}), 400

        # Validate message length
        if len(data['message']) > 2000:
            return jsonify({'error': 'Message too long (max 2000 characters)'}), 400

        # Log feedback to file
        import datetime
        log_entry = f"""
{'='*60}
Timestamp: {data.get('timestamp', datetime.datetime.now().isoformat())}
Name: {data.get('name', 'Anonymous')}
Email: {data.get('email', 'N/A')}
Type: {data.get('type', 'N/A')}
URL: {data.get('url', 'N/A')}
Message:
{data['message']}
{'='*60}
"""

        # Write to feedback log file
        with open('feedback.log', 'a', encoding='utf-8') as f:
            f.write(log_entry)

        return jsonify({'success': True, 'message': 'Feedback received'}), 200
    except Exception as e:
        return jsonify({'error': f'Error submitting feedback: {str(e)}'}), 500

if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=5000)
